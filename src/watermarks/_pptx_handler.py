"""
PPTX 格式水印处理：嵌入/提取零宽字符。

仅负责文件格式操作（读写 PPTX 文本框内容），
编解码逻辑由 office_wm.py 统一处理。
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from loguru import logger


def embed_pptx(
    input_path: Path, zwc_block: str, output_path: Path,
) -> tuple[bool, str]:
    """
    在 PPTX 文件中嵌入零宽字符块。

    遍历所有 slide → shape → paragraph → run，
    在第一个非空 run 的文本前端插入 ZWC 块。
    必须检查 shape.has_text_frame 避免 AttributeError。

    Returns:
        (success, message) 元组
    """
    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        return False, f"Cannot open PPTX: {e}"

    inserted = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip():
                        run.text = zwc_block + run.text
                        inserted = True
                        break
                if inserted:
                    break
            if inserted:
                break
        if inserted:
            break

    if not inserted:
        return False, "No text content found in PPTX"

    try:
        prs.save(str(output_path))
    except Exception as e:
        return False, f"Cannot save PPTX: {e}"

    logger.debug(f"ZWC block inserted into PPTX: {input_path.name}")
    return True, "OK"


def extract_pptx(file_path: Path) -> Optional[str]:
    """
    从 PPTX 文件提取所有文本内容。

    遍历所有 slide → shape → paragraph → run，
    拼接全部文本供 zwc_decode 搜索水印标记。

    Returns:
        拼接后的全文文本，失败返回 None
    """
    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        logger.warning(f"Cannot open PPTX for extraction: {e}")
        return None

    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        parts.append(run.text)
    return "".join(parts) if parts else None
