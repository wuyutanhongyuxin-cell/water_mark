"""
测试共享 fixtures。

提供固定密钥、标准载荷、各类型测试文件的程序化生成。
所有 fixture 使用 tmp_path 确保测试隔离。
"""

import json
import os
from pathlib import Path

import numpy as np
import pytest

from src.watermarks.base import WatermarkPayload


# ========== 固定密钥（全局 autouse） ==========

@pytest.fixture(autouse=True)
def fixed_key(monkeypatch):
    """设置固定密钥环境变量，确保测试确定性。"""
    monkeypatch.setenv("WATERMARK_MASTER_KEY", "a" * 64)


@pytest.fixture(autouse=True)
def disable_ai(monkeypatch):
    """确保测试中 AI 功能始终禁用。"""
    monkeypatch.delenv("DEEPSEEK_API_KEY", raising=False)


@pytest.fixture(autouse=True)
def reset_router_cache():
    """每次测试前清除路由缓存（lru_cache）。"""
    from src.core.router import _load_rules_cached, _load_settings_cached
    _load_rules_cached.cache_clear()
    _load_settings_cached.cache_clear()
    yield
    _load_rules_cached.cache_clear()
    _load_settings_cached.cache_clear()


@pytest.fixture(autouse=True)
def reset_audit_state():
    """重置审计日志初始化状态，避免测试间互相影响。"""
    import src.security.audit as audit_mod
    audit_mod._initialized = False
    audit_mod._audit_logger = None


# ========== 标准载荷 ==========

@pytest.fixture
def sample_payload() -> WatermarkPayload:
    """标准测试载荷。"""
    return WatermarkPayload(
        employee_id="E001",
        timestamp="2026-04-05T00:00:00Z",
        file_hash="abcdef1234567890",
    )


# ========== 图像生成 ==========

@pytest.fixture
def sample_image(tmp_path) -> Path:
    """生成 300x300 RGB 随机噪声 PNG 图像。"""
    import cv2
    rng = np.random.RandomState(42)
    img = rng.randint(0, 256, (300, 300, 3), dtype=np.uint8)
    path = tmp_path / "test_image.png"
    _, buf = cv2.imencode(".png", img)
    path.write_bytes(buf.tobytes())
    return path


@pytest.fixture
def small_image(tmp_path) -> Path:
    """生成 100x100 小图像（用于边界测试）。"""
    import cv2
    rng = np.random.RandomState(99)
    img = rng.randint(0, 256, (100, 100, 3), dtype=np.uint8)
    path = tmp_path / "small_image.png"
    _, buf = cv2.imencode(".png", img)
    path.write_bytes(buf.tobytes())
    return path


# ========== 文本生成 ==========

@pytest.fixture
def sample_txt(tmp_path) -> Path:
    """生成纯文本文件。"""
    path = tmp_path / "test.txt"
    path.write_text("Hello World\nThis is a test file.\nLine three.\n",
                    encoding="utf-8")
    return path


@pytest.fixture
def sample_csv(tmp_path) -> Path:
    """生成 CSV 文件。"""
    path = tmp_path / "test.csv"
    path.write_text("name,age,dept\nAlice,30,Engineering\nBob,25,Marketing\n",
                    encoding="utf-8")
    return path


@pytest.fixture
def sample_json(tmp_path) -> Path:
    """生成 JSON 文件。"""
    path = tmp_path / "test.json"
    data = {"name": "Alice", "dept": "Engineering", "level": 3}
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")
    return path


@pytest.fixture
def sample_md(tmp_path) -> Path:
    """生成 Markdown 文件。"""
    path = tmp_path / "test.md"
    path.write_text("# Test\n\nThis is a **test** document.\n\n- Item 1\n- Item 2\n",
                    encoding="utf-8")
    return path


# ========== PDF 生成 ==========

@pytest.fixture
def sample_pdf(tmp_path) -> Path:
    """用 PyMuPDF 创建含文本的单页 PDF。"""
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "This is a test PDF document.", fontsize=14)
    page.insert_text((72, 120), "Second line of text for testing.", fontsize=12)
    path = tmp_path / "test.pdf"
    doc.save(str(path))
    doc.close()
    return path


# ========== Office 生成 ==========

@pytest.fixture
def sample_docx(tmp_path) -> Path:
    """生成 DOCX 文件。"""
    import docx
    document = docx.Document()
    document.add_paragraph("This is a test DOCX document.")
    document.add_paragraph("Second paragraph for testing.")
    path = tmp_path / "test.docx"
    document.save(str(path))
    return path


@pytest.fixture
def sample_xlsx(tmp_path) -> Path:
    """生成 XLSX 文件。"""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Department"
    ws["A2"] = "Alice"
    ws["B2"] = "Engineering"
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    wb.close()
    return path


@pytest.fixture
def sample_pptx(tmp_path) -> Path:
    """生成 PPTX 文件。"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "This is test content."
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


# ========== 音频生成 ==========

@pytest.fixture
def sample_wav(tmp_path) -> Path:
    """生成 44100Hz 正弦波 WAV 文件（≥65536 samples）。"""
    import soundfile as sf
    sr = 44100
    duration = 2.0  # 2 秒 = 88200 samples > 65536
    t = np.linspace(0, duration, int(sr * duration), endpoint=False)
    signal = 0.5 * np.sin(2 * np.pi * 440 * t)  # 440Hz A4
    path = tmp_path / "test.wav"
    sf.write(str(path), signal, sr, format="WAV", subtype="PCM_16")
    return path


# ========== 视频生成 ==========

@pytest.fixture
def sample_avi(tmp_path) -> Path:
    """生成 20 帧 320x320 随机噪声 AVI 视频。"""
    import cv2
    path = tmp_path / "test.avi"
    fourcc = cv2.VideoWriter_fourcc(*"FFV1")
    writer = cv2.VideoWriter(str(path), fourcc, 10.0, (320, 320))
    rng = np.random.RandomState(42)
    for _ in range(20):
        frame = rng.randint(0, 256, (320, 320, 3), dtype=np.uint8)
        writer.write(frame)
    writer.release()
    return path
